import os
import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
import numpy_financial as npf  # Only for IRR
from datetime import datetime

# -------------------------
# Data directory
# -------------------------
DATA_DIR = os.path.join(os.path.dirname(__file__), "data")

# -------------------------
# Helpers
# -------------------------
def _get_series(df):
    if df is None or df.empty:
        return None, None
    col_time = df.columns[0]
    col_val = df.columns[1]
    try:
        t = pd.to_datetime(df[col_time])
    except Exception:
        t = df[col_time]
    v = pd.to_numeric(df[col_val], errors="coerce").fillna(0.0)
    return t, v

def _payback_year(cashflows):
    cum = np.cumsum(cashflows)
    for i, c in enumerate(cum, start=1):
        if c >= 0:
            return i
    return None

def fmt(x):
    if isinstance(x, (int, float)):
        return f"{x:,.2f}"
    return x

# -------------------------
# Load benchmark and solar profiles
# -------------------------
def load_profiles():
    profiles = {}
    missing_profiles = []

    def try_load(name, filename):
        path = os.path.join(DATA_DIR, filename)
        if os.path.exists(path):
            profiles[name] = pd.read_csv(path)
        else:
            profiles[name] = None
            missing_profiles.append(filename)  # Track missing files

    # ---- Benchmark profiles ----
    try_load("Benchmark_Office", "Benchmark_Profile_Office.csv")
    try_load("Benchmark_Storage", "Benchmark_Profile_Storage.csv")
    try_load("Benchmark_RetailPark", "Benchmark_Profile_RetailPark.csv")
    try_load("Benchmark_Logistics", "Benchmark_Profile_Logistics.csv")
    try_load("Benchmark_LightManufacturing", "Benchmark_Profile_LightManufacturing.csv")
    try_load("Benchmark_HeavyManufacturing", "Benchmark_Profile_HeavyManufacturing.csv")

    # ---- Solar profiles ----
    try_load("Solar_South", "Solar_Profile_South.csv")
    try_load("Solar_Midlands", "Solar_Profile_Midlands.csv")
    try_load("Solar_Scotland", "Solar_Profile_Scotland.csv")

    return profiles, missing_profiles


# -------------------------
# Energy + Financial calculations
# -------------------------
def prepare_year1_timeseries(system_size_kwp, solar_df, demand_df,
                             used_regional_solar=False, used_benchmark_demand=False,
                             site_type="Office", floor_space_m2=0.0,
                             intensity_office=65.0, intensity_storage=27.0,
                             export_allowed=True):
    t_s, v_s = _get_series(solar_df)
    if used_regional_solar:
        solar_hh = (v_s.astype(float) * float(system_size_kwp)).values
    else:
        solar_hh = v_s.astype(float).values

    t_d, v_d = _get_series(demand_df)
    if used_benchmark_demand:
        intensity = intensity_office if site_type == "Office" else intensity_storage
        annual_demand = float(intensity) * float(floor_space_m2)
        shape = v_d.astype(float).values
        shape_sum = shape.sum()
        demand_hh = np.zeros_like(shape) if shape_sum <= 0 else shape / shape_sum * annual_demand
    else:
        demand_hh = v_d.astype(float).values

    n = min(len(solar_hh), len(demand_hh))
    solar_hh = solar_hh[:n]
    demand_hh = demand_hh[:n]
    if t_s is not None and len(t_s) >= n:
        index = t_s[:n]
    else:
        index = pd.RangeIndex(0, n, 1)

    return index, solar_hh, demand_hh

def aggregate_monthly_savings(index, solar_hh, demand_hh, model,
                              import_tariff, export_tariff, ppa_rate,
                              export_allowed=True):
    df = pd.DataFrame({
        "ts": pd.to_datetime(index, errors="coerce"),
        "solar": solar_hh,
        "demand": demand_hh
    })
    if df["ts"].isna().all():
        df["month"] = ((df.index.values / (len(df) / 12))).astype(int) + 1
        df["month"] = df["month"].clip(1, 12)
    else:
        df["month"] = df["ts"].dt.month

    self_cons = np.minimum(df["solar"], df["demand"])
    export = np.maximum(df["solar"] - df["demand"], 0.0) if export_allowed else 0.0

    if model == "Owner Occupier":
        savings_gbp = self_cons * import_tariff + export * export_tariff
    else:
        savings_gbp = self_cons * max(import_tariff - ppa_rate, 0.0)

    monthly = savings_gbp.groupby(df["month"]).sum()
    monthly = monthly.reindex(range(1, 13), fill_value=0.0)
    return monthly

def calculate_project_financials(model, system_size_kwp, solar_df, demand_df,
                                 project_life=25, capex_per_kwp=800, opex_per_kwp=15,
                                 import_tariff=0.25, export_tariff=0.08, ppa_rate=0.18,
                                 replace_years=[15], inflation=0.0,
                                 export_allowed=True, used_regional_solar=False,
                                 used_benchmark_demand=False, site_type="Office",
                                 floor_space_m2=0.0,
                                 inverter_replacement_cost_per_kwp=50.0):
    degradation = 0.005

    index, solar_hh, demand_hh = prepare_year1_timeseries(
        system_size_kwp, solar_df, demand_df,
        used_regional_solar=used_regional_solar,
        used_benchmark_demand=used_benchmark_demand,
        site_type=site_type, floor_space_m2=floor_space_m2,
        export_allowed=export_allowed
    )

    annual_yield_y1 = float(np.sum(solar_hh))
    self_consumption_y1 = float(np.sum(np.minimum(solar_hh, demand_hh)))
    export_y1 = float(np.sum(np.maximum(solar_hh - demand_hh, 0.0))) if export_allowed else 0.0
    cons_ratio = (self_consumption_y1 / annual_yield_y1) if annual_yield_y1 > 0 else 0.0

    years = list(range(1, project_life + 1))
    gen_y = [annual_yield_y1 * ((1 - degradation) ** (y - 1)) for y in years]
    selfc_y = [g * cons_ratio for g in gen_y]
    export_y = [(g - s) if export_allowed else 0.0 for g, s in zip(gen_y, selfc_y)]

    capex = float(system_size_kwp) * float(capex_per_kwp)
    annual_opex_y = [float(system_size_kwp) * float(opex_per_kwp) * ((1 + float(inflation)) ** (y - 1)) for y in years]

    owner_cf, landlord_cf, tenant_cf = [], [], []

    if model == "Owner Occupier":
        for i, y in enumerate(years):
            inflow = selfc_y[i] * import_tariff + export_y[i] * export_tariff
            net = inflow - annual_opex_y[i]
            if y == 1:
                net -= capex
            if y in replace_years:
                net -= float(inverter_replacement_cost_per_kwp) * float(system_size_kwp)
            owner_cf.append(net)
    else:
        for i, y in enumerate(years):
            landlord_revenue = selfc_y[i] * ppa_rate + export_y[i] * export_tariff
            net_landlord = landlord_revenue - annual_opex_y[i]
            if y == 1:
                net_landlord -= capex
            if y in replace_years:
                net_landlord -= float(inverter_replacement_cost_per_kwp) * float(system_size_kwp)
            landlord_cf.append(net_landlord)

            tenant_sav = selfc_y[i] * max(import_tariff - ppa_rate, 0.0)
            tenant_cf.append(tenant_sav)

    def _safe_irr(cf):
        try:
            val = float(npf.irr(cf))
            if np.isfinite(val):
                return val
            return None
        except Exception:
            return None

    if model == "Owner Occupier":
        irr_owner = _safe_irr(owner_cf)
        pb_owner = _payback_year(owner_cf)
        gross_inflow_life = sum([selfc_y[i]*import_tariff + export_y[i]*export_tariff for i in range(len(years))])
        replacements_total = sum(float(inverter_replacement_cost_per_kwp) * float(system_size_kwp) for y in years if y in replace_years)
        net_lifetime_value = gross_inflow_life - sum(annual_opex_y) - capex - replacements_total
        results = {
            "cashflow_yearly": pd.DataFrame({"Year": years, "Owner": owner_cf}),
            "irr": {"Owner Occupier": irr_owner},
            "payback": {"Owner Occupier": pb_owner},
            "lifetime": {"Owner Occupier": {
                "Capex": capex,
                "Annual Opex (Y1)": annual_opex_y[0],
                "Net lifetime savings": net_lifetime_value,
                "Net lifetime income": net_lifetime_value,
            }},
        }
    else:
        irr_landlord = _safe_irr(landlord_cf)
        pb_landlord = _payback_year(landlord_cf)
        landlord_net_lifetime_income = sum(landlord_cf)
        tenant_net_lifetime_savings = sum(tenant_cf)
        results = {
            "cashflow_yearly": pd.DataFrame({"Year": years, "Landlord": landlord_cf}),
            "irr": {"Landlord": irr_landlord},
            "payback": {"Landlord": pb_landlord},
            "lifetime": {
                "Landlord": {
                    "Capex": capex,
                    "Annual Opex (Y1)": annual_opex_y[0],
                    "Net lifetime income": landlord_net_lifetime_income,
                },
                "Tenant": {
                    "Net lifetime savings": tenant_net_lifetime_savings
                }
            },
        }

    technical_y1 = {
        "Annual Yield (kWh)": annual_yield_y1,
        "Consumed on site (kWh)": self_consumption_y1,
        "Exported (kWh)": export_y1,
    }

    monthly_savings_y1 = aggregate_monthly_savings(
        index, solar_hh, demand_hh, model, import_tariff, export_tariff, ppa_rate, export_allowed=export_allowed
    )

    return results, technical_y1, monthly_savings_y1

# -------------------------
# PPT export
# -------------------------

def export_ppt(template_bytes, project_name, summary, financials_view, technical_y1,
               monthly_chart_buf, cashflow_chart_buf, monthly_title, cashflow_title,
               layout):
    """
    Create a PPTX from a template (BytesIO) and insert tables/charts according to layout.
    layout keys: 'summary', 'technical', 'financial', 'monthly', 'cashflow'
    Each layout element contains: slide, x, y, w, h (inches)
    """

    # Load the presentation (uploaded template or blank)
    if template_bytes:
        try:
            prs = Presentation(template_bytes)
        except Exception:
            prs = Presentation()
    else:
        prs = Presentation()

    # Helper: get slide by number (auto-create if not enough slides)
    def get_slide(num):
        idx = max(0, num - 1)
        while len(prs.slides) <= idx:
            prs.slides.add_slide(prs.slide_layouts[6])  # blank
        return prs.slides[idx]

    # ------------------------
    # Add SUMMARY TABLE
    # ------------------------
    try:
        cfg = layout["summary"]
        slide = get_slide(cfg["slide"])
        left, top = Inches(cfg["x"]), Inches(cfg["y"])
        width, height = Inches(cfg["w"]), Inches(cfg["h"])

        rows = len(summary) + 1
        cols = 2
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        table.cell(0, 0).text = "Parameter"
        table.cell(0, 1).text = "Value"

        r = 1
        for k, v in summary.items():
            table.cell(r, 0).text = str(k)
            table.cell(r, 1).text = str(v)
            r += 1
    except Exception:
        pass

    # ------------------------
    # Add TECHNICAL TABLE
    # ------------------------
    try:
        cfg = layout["technical"]
        slide = get_slide(cfg["slide"])
        left, top = Inches(cfg["x"]), Inches(cfg["y"])
        width, height = Inches(cfg["w"]), Inches(cfg["h"])

        table = slide.shapes.add_table(4, 2, left, top, width, height).table

        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"

        table.cell(1, 0).text = "Annual Yield (kWh)"
        table.cell(1, 1).text = f"{technical_y1['Annual Yield (kWh)']:,.0f}"

        table.cell(2, 0).text = "Consumed on site (kWh)"
        table.cell(2, 1).text = f"{technical_y1['Consumed on site (kWh)']:,.0f}"

        table.cell(3, 0).text = "Exported (kWh)"
        table.cell(3, 1).text = f"{technical_y1['Exported (kWh)']:,.0f}"
    except Exception:
        pass

    # ------------------------
    # Add FINANCIAL TABLE
    # ------------------------
    try:
        cfg = layout["financial"]
        slide = get_slide(cfg["slide"])
        left, top = Inches(cfg["x"]), Inches(cfg["y"])
        width, height = Inches(cfg["w"]), Inches(cfg["h"])

        rows = len(financials_view)
        cols = len(financials_view[0])

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        for i, row in enumerate(financials_view):
            for j, val in enumerate(row):
                table.cell(i, j).text = str(val)
    except Exception:
        pass

    # ------------------------
    # Add MONTHLY CHART
    # ------------------------
    try:
        cfg = layout["monthly"]
        slide = get_slide(cfg["slide"])
        left, top = Inches(cfg["x"]), Inches(cfg["y"])
        width, height = Inches(cfg["w"]), Inches(cfg["h"])

        slide.shapes.add_picture(monthly_chart_buf, left, top, width=width, height=height)
    except Exception:
        pass

    # ------------------------
    # Add CASHFLOW CHART
    # ------------------------
    try:
        cfg = layout["cashflow"]
        slide = get_slide(cfg["slide"])
        left, top = Inches(cfg["x"]), Inches(cfg["y"])
        width, height = Inches(cfg["w"]), Inches(cfg["h"])

        slide.shapes.add_picture(cashflow_chart_buf, left, top, width=width, height=height)
    except Exception:
        pass

    # Return BytesIO buffer
    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# -------------------------
# Streamlit app
# -------------------------

def main():
    st.set_page_config(page_title="Solar Modelling Tool")

    col1, col2 = st.columns([6, 1])
    with col1:
        st.title("☀️ Solar Modelling Tool")
    with col2:
        pass
    profiles, missing = load_profiles()

    # --- Sidebar ---
    with st.sidebar.expander("🛠 Debug Tools", expanded=False):
    
        st.write("**DATA_DIR:**", DATA_DIR)
    
        # File visibility
        try:
            st.write("**Files detected:**", os.listdir(DATA_DIR))
        except Exception as e:
            st.error(f"Error reading DATA_DIR: {e}")
    
        if missing:
            st.warning("⚠️ Missing profiles:\n" + "\n".join(f"- {m}" for m in missing))
        else:
            st.success("✅ All profile files successfully loaded.")
    
    with st.sidebar.expander("📈 Savills CAPEX Curve Parameters", expanded=False):
    
        a_value = st.number_input("Coefficient a", value=1398.58238)
        b_value = st.number_input("Exponent b", value=-0.10814)
    
        st.caption("Formula: **CAPEX (£/kWp) = a × (System Size)^b**")




    # -------------------------
    # PowerPoint layout settings (user editable)
    # -------------------------
    with st.sidebar.expander("📐 PPT Layout Settings", expanded=False):
    
        layout = {}
    
        layout["summary"] = {
            "slide": st.number_input("Summary slide", 1, 20, 1),
            "x": st.number_input("Summary X (inches)", 0.0, 20.0, 1.0),
            "y": st.number_input("Summary Y (inches)", 0.0, 20.0, 1.5),
            "w": st.number_input("Summary Width (inches)", 1.0, 20.0, 8.0),
            "h": st.number_input("Summary Height (inches)", 1.0, 20.0, 3.0),
        }
    
        layout["technical"] = {
            "slide": st.number_input("Technical slide", 1, 20, 2),
            "x": st.number_input("Technical X (inches)", 0.0, 20.0, 1.0),
            "y": st.number_input("Technical Y (inches)", 0.0, 20.0, 1.5),
            "w": st.number_input("Technical Width (inches)", 1.0, 20.0, 8.0),
            "h": st.number_input("Technical Height (inches)", 1.0, 20.0, 3.0),
        }
    
        layout["financial"] = {
            "slide": st.number_input("Financial slide", 1, 20, 3),
            "x": st.number_input("Financial X (inches)", 0.0, 20.0, 1.0),
            "y": st.number_input("Financial Y (inches)", 0.0, 20.0, 1.5),
            "w": st.number_input("Financial Width (inches)", 1.0, 20.0, 8.0),
            "h": st.number_input("Financial Height (inches)", 1.0, 20.0, 3.0),
        }
    
        layout["monthly"] = {
            "slide": st.number_input("Monthly chart slide", 1, 20, 4),
            "x": st.number_input("Monthly X (inches)", 0.0, 20.0, 1.0),
            "y": st.number_input("Monthly Y (inches)", 0.0, 20.0, 1.0),
            "w": st.number_input("Monthly Width (inches)", 1.0, 20.0, 8.0),
            "h": st.number_input("Monthly Height (inches)", 1.0, 20.0, 4.0),
        }
    
        layout["cashflow"] = {
            "slide": st.number_input("Cashflow slide", 1, 20, 5),
            "x": st.number_input("Cashflow X (inches)", 0.0, 20.0, 1.0),
            "y": st.number_input("Cashflow Y (inches)", 0.0, 20.0, 1.0),
            "w": st.number_input("Cashflow Width (inches)", 1.0, 20.0, 8.0),
            "h": st.number_input("Cashflow Height (inches)", 1.0, 20.0, 4.0),
        }
    # Initialise session state
    for k in ["results", "technical_y1", "monthly_savings_y1",
              "summary_dict", "fin_view", "project_name",
              "monthly_title", "cashflow_title", "ppt_buf"]:
        st.session_state.setdefault(k, None)

    # --- Inputs ---
    st.subheader("Project Name")
    project_name = st.text_input("Enter a name for the project",
                                 st.session_state.get("project_name") or "Savills Solar Project")

    st.subheader("Demand Profile")
    demand_option = st.radio("Do you have half-hourly demand profile?",
                             ["No - Use Benchmark Profile","Yes - Upload CSV"],index=0)
    
    floor_space_m2 = 0.0
    site_type = "Office"
    if demand_option == "Yes - Upload CSV":
        demand_file = st.file_uploader("Upload demand CSV", type="csv", key="demand_upload")
        demand_profile = pd.read_csv(demand_file) if demand_file else None
        used_benchmark_demand = False
    else:
        site_type = st.selectbox(
            "Select Site Type",
            ["Office", "Storage", "Retail Park", "Logistics", "Light Manufacturing", "Heavy Manufacturing"],
            help="Choose the category that best represents the site's energy use pattern. \
        If unsure, select the closest operational type."
        )

        floor_space_m2 = st.number_input("Floor space (m²)", min_value=0.0, value=1000.0, step=10.0)
        benchmark_map = {
            "Office": "Benchmark_Office",
            "Storage": "Benchmark_Storage",
            "Retail Park": "Benchmark_RetailPark",
            "Logistics": "Benchmark_Logistics",
            "Light Manufacturing": "Benchmark_LightManufacturing",
            "Heavy Manufacturing": "Benchmark_HeavyManufacturing",
        }
        benchmark_key = benchmark_map[site_type]
        # 🔍 SAFETY CHECK — file missing?
        if profiles[benchmark_key] is None:
            st.error(f"The benchmark profile for **{site_type}** has not been uploaded yet.\n\n"
                    "Please upload the CSV file into the `data/` folder.")
            st.stop()
        demand_profile = profiles[benchmark_key]
        used_benchmark_demand = True

    st.subheader("Solar Profile")
    solar_option = st.radio(
        "Do you have half-hourly solar profile?",
        ["No - Use Regional Profile", "Yes - Upload CSV"],
        index=0,
        help="Upload your own site-specific solar generation data if available. \
    Otherwise use the pre-loaded regional profiles."
    )

    if solar_option == "Yes - Upload CSV":
        solar_file = st.file_uploader("Upload solar CSV", type="csv", key="solar_upload")
        solar_profile = pd.read_csv(solar_file) if solar_file else None
        used_regional_solar = False
    else:
        region = st.selectbox(
            "Select Region",
            ["South", "Midlands", "Scotland"],
            help="Select the region where the site is located. \
        This chooses the appropriate regional solar irradiance profile."
        )
        solar_profile = profiles[f"Solar_{region}"]
        used_regional_solar = True

    st.subheader("System & Financial Inputs")
    system_size = st.number_input("System Size (kWp)", min_value=10.0, max_value=500000.0, value=500.0, step=10.0)
    capex_option = st.radio(
        "CAPEX Method",
        ["Savills Database", "Direct CAPEX Input"],
        index=0,   # default = first option
        help="Select whether to use the Savills capex curve or manually enter CAPEX (£/kWp)."
    )
    if capex_option == "Direct CAPEX Input":
        capex_per_kw = st.number_input("CAPEX (£/kWp)", 0.0, 5000.0, 800.0)
    else:
        # Use Savills database formula: y = a * x^b
        capex_per_kw = a_value * (system_size ** b_value)
        st.info(f"Calculated CAPEX from Savills DB: **£{capex_per_kw:,.0f}/kWp**")
    opex_per_kw = st.number_input("O&M (£/kWp/year)", 0.0, 200.0, 15.0)
    inverter_replacement_cost_per_kwp = st.number_input("Inverter Replacement Cost (£/kWp)", 0.0, 1000.0, 50.0, step=1.0)
    model = st.radio("Financial Model", ["Owner Occupier", "Landlord Funded (PPA to Tenant)"])
    import_tariff = st.number_input("Import Tariff (£/kWh)", 0.0, 1.0, 0.25)
    export_tariff = st.number_input("Export Tariff (£/kWh)", 0.0, 1.0, 0.08)
    ppa_rate = st.number_input("PPA Rate (£/kWh)", 0.0, 1.0, 0.18) if model == "Landlord Funded (PPA to Tenant)" else None
    project_life = st.number_input("Project Lifespan (years)", 1, 50, 25)
    inflation = st.number_input("Annual Inflation Rate (%)", 0.0, 10.0, 0.0) / 100
    replace_years = st.multiselect("Inverter Replacement Years", list(range(1, 51)), [15])
    export_allowed = st.checkbox("Export Allowed?", True)

    # --- Run simulation: compute + store, don't render here ---
    if st.button("Run / Update Simulation", type="primary"):
        if demand_profile is not None and solar_profile is not None:
            results, technical_y1, monthly_savings_y1 = calculate_project_financials(
                model, system_size, solar_profile, demand_profile,
                project_life, capex_per_kw, opex_per_kw,
                import_tariff, export_tariff, ppa_rate or 0.0,
                replace_years, inflation, export_allowed,
                used_regional_solar, used_benchmark_demand,
                site_type, floor_space_m2, inverter_replacement_cost_per_kwp
            )

            summary_dict = {
                "Project Name": project_name,
                "System Size (kWp)": fmt(system_size),
                "Financial Model": model,
                "CAPEX (£/kWp)": fmt(capex_per_kw),
                "OPEX (£/kWp/year)": fmt(opex_per_kw),
                "Inverter Replacement (£/kWp)": fmt(inverter_replacement_cost_per_kwp),
                "Import Tariff (£/kWh)": fmt(import_tariff),
                "Export Tariff (£/kWh)": fmt(export_tariff),
                "PPA Rate (£/kWh)": fmt(ppa_rate) if ppa_rate else "-",
                "Project Lifespan (years)": project_life,
                "Inverter Replacement Years": ", ".join(map(str, replace_years)),
                "Inflation Rate": f"{inflation*100:.2f}%",
                "Export Allowed": "Yes" if export_allowed else "No",
                "Demand Source": "Benchmark" if used_benchmark_demand else "Uploaded",
                "Site Type": site_type if used_benchmark_demand else "-",
                "Floor space (m²)": fmt(floor_space_m2) if used_benchmark_demand else "-",
                "Solar Source": "Regional" if used_regional_solar else "Uploaded",
                "Degradation": "0.5%/yr",
            }

            # Build FINANCIAL METRICS view (compact, not 25 rows)
            fin_view = []
            if model == "Owner Occupier":
                headers = ["Metric", "Owner Occupier"]
                fin_view.append(headers)
                capex = results["lifetime"]["Owner Occupier"]["Capex"]
                opex_y1 = results["lifetime"]["Owner Occupier"]["Annual Opex (Y1)"]
                net_lifetime_savings = results["lifetime"]["Owner Occupier"]["Net lifetime savings"]
                net_lifetime_income = results["lifetime"]["Owner Occupier"]["Net lifetime income"]
                irr = results["irr"]["Owner Occupier"]
                pb = results["payback"]["Owner Occupier"]
                rows = [
                    ["Capex", f"£{capex:,.0f}"],
                    ["Annual Opex", f"£{opex_y1:,.0f} (Y1)"],
                    ["Net lifetime savings", f"£{net_lifetime_savings:,.0f}"],
                    ["Net lifetime income", f"£{net_lifetime_income:,.0f}"],
                    ["IRR", f"{irr*100:.1f}%" if irr else "N/A"],
                    ["Payback years", f"{pb}" if pb else "N/A"],
                ]
                fin_view.extend(rows)
            else:
                headers = ["Metric", "Landlord", "Tenant"]
                fin_view.append(headers)
                capex = results["lifetime"]["Landlord"]["Capex"]
                opex_y1 = results["lifetime"]["Landlord"]["Annual Opex (Y1)"]
                landlord_net_income = results["lifetime"]["Landlord"]["Net lifetime income"]
                tenant_net_savings = results["lifetime"]["Tenant"]["Net lifetime savings"]
                irr_l = results["irr"]["Landlord"]
                pb_l = results["payback"]["Landlord"]
                rows = [
                    ["Capex", f"£{capex:,.0f}", "N/A"],
                    ["Annual Opex", f"£{opex_y1:,.0f} (Y1)", "N/A"],
                    ["Net lifetime savings", "N/A", f"£{tenant_net_savings:,.0f}"],
                    ["Net lifetime income", f"£{landlord_net_income:,.0f}", "N/A"],
                    ["IRR", f"{irr_l*100:.1f}%" if irr_l else "N/A", "N/A"],
                    ["Payback years", f"{pb_l}" if pb_l else "N/A", "N/A"],
                ]
                fin_view.extend(rows)

            # --- Charts for PDF buffers ---
            months = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
            monthly_title = f"Monthly Savings (Year 1) ({'Owner Occupier' if model=='Owner Occupier' else 'Tenant'})"
            fig_m, ax_m = plt.subplots()
            ax_m.bar(months, monthly_savings_y1.values)
            ax_m.set_ylabel("£")
            ax_m.set_xlabel("Month")
            fig_m.tight_layout()
            m_buf = BytesIO()
            fig_m.savefig(m_buf, format="png"); m_buf.seek(0)

            cashflow_for = "Owner Occupier" if model == "Owner Occupier" else "Landlord"
            cashflow_title = f"Cashflow ({cashflow_for})"
            cash_df = results["cashflow_yearly"]
            fig_c, ax_c = plt.subplots()
            yearly_cf = cash_df["Owner"].values if model == "Owner Occupier" else cash_df["Landlord"].values
            years = cash_df["Year"].values
            ax_c.plot(years, yearly_cf, marker="o", label="Annual cashflow")
            cum_cf = np.cumsum(yearly_cf)
            ax_c.plot(years, cum_cf, marker="o", linestyle="--", label="Cumulative cashflow")
            ax_c.legend()
            fig_c.tight_layout()
            c_buf = BytesIO()
            fig_c.savefig(c_buf, format="png"); c_buf.seek(0)

            # PDF uses the compact fin_view
            ppt_template_file = st.file_uploader("Upload PowerPoint Template (optional)", type=["pptx"], key="ppt_template")

            # Build PPT report (use uploaded template if provided)
            ppt_buf = export_ppt(
                ppt_template_file,
                project_name,
                summary_dict,
                [fin_view[0]] + fin_view[1:],
                technical_y1,
                m_buf,
                c_buf,
                monthly_title,
                cashflow_title,
                layout
            )

            # Save everything in session_state
            st.session_state.update({
                "results": results,
                "technical_y1": technical_y1,
                "monthly_savings_y1": monthly_savings_y1,
                "summary_dict": summary_dict,
                "fin_view": fin_view,
                "project_name": project_name,
                "monthly_title": monthly_title,
                "cashflow_title": cashflow_title,
                "ppt_buf": ppt_buf
            })

    # --- Always render results if available ---
    if st.session_state.get("results"):
        results = st.session_state["results"]
        technical_y1 = st.session_state["technical_y1"]
        summary_dict = st.session_state["summary_dict"]
        monthly_savings_y1 = st.session_state["monthly_savings_y1"]
        fin_view = st.session_state["fin_view"]

        st.subheader("Summary Inputs")
        st.dataframe(pd.DataFrame(list(summary_dict.items()), columns=["Parameter", "Value"]),
                     hide_index=True, use_container_width=True)

        st.subheader("Technical Output (Year 1)")
        tech_df = pd.DataFrame({
            "Metric": ["Annual Yield (kWh)", "Consumed on site (kWh)", "Exported (kWh)"],
            "Value": [f"{technical_y1['Annual Yield (kWh)']:,.0f}",
                      f"{technical_y1['Consumed on site (kWh)']:,.0f}",
                      f"{technical_y1['Exported (kWh)']:,.0f}"]
        })
        st.dataframe(tech_df, hide_index=True, use_container_width=True)

        # --- Correct Financial Metrics table with matching header background ---
        st.subheader("Financial Metrics")
        headers = fin_view[0]
        rows = fin_view[1:]
        fin_df = pd.DataFrame(rows, columns=headers).set_index("Metric")
        styled_fin_df = fin_df.style.set_table_styles(
            [{"selector": "thead th", "props": [("background-color", "#f0f0f0"),
                                                ("font-weight", "bold")]}]
        )
        st.dataframe(styled_fin_df, use_container_width=True)

        # --- Charts (re-render from state) ---
        months = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
        st.subheader(st.session_state["monthly_title"])
        fig_m2, ax_m2 = plt.subplots()
        ax_m2.bar(months, monthly_savings_y1.values)
        ax_m2.set_ylabel("£")
        ax_m2.set_xlabel("Month")
        st.pyplot(fig_m2)

        st.subheader(st.session_state["cashflow_title"])
        fig_c2, ax_c2 = plt.subplots()
        cash_df = results["cashflow_yearly"]
        yearly_cf = cash_df["Owner"].values if "Owner" in cash_df.columns else cash_df["Landlord"].values
        years = cash_df["Year"].values
        ax_c2.plot(years, yearly_cf, marker="o", label="Annual cashflow")
        cum_cf = np.cumsum(yearly_cf)
        ax_c2.plot(years, cum_cf, marker="o", linestyle="--", label="Cumulative cashflow")
        ax_c2.legend()
        st.pyplot(fig_c2)

    # --- Persistent Download Button (single instance) ---
    if st.session_state.get("ppt_buf"):
        # Format timestamp for filename
        timestamp = datetime.now().strftime("%y%m%d_%H%M")
        safe_project_name = st.session_state["project_name"].replace(" ", "_")  # avoid spaces in filename
        file_name = f"{timestamp}_{safe_project_name}.pptx"
    
        st.download_button("Download PPT Report",
                           data=st.session_state["ppt_buf"],
                           file_name=file_name,
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                           key="download_ppt_persist")


if __name__ == "__main__":
    main()










